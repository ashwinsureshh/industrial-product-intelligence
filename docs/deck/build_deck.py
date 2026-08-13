"""Fill the UniHack prototype template with the project's measured results.

Every number here is reproducible from the repo:
  run_benchmark.py · run_hybrid.py · run_delivery_accuracy.py · run_cost_model.py

The template's slide order, count and branding are preserved. Only slide 1 is
repurposed (organizers' guidelines -> cover), per the user's decision.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- brand, sampled from the template's own background art -------------------
NAVY   = RGBColor(0x00, 0x38, 0x78)
BLUE   = RGBColor(0x02, 0x96, 0xFC)
INK    = RGBColor(0x16, 0x20, 0x2C)
MUTED  = RGBColor(0x5B, 0x6B, 0x7C)
CARD   = RGBColor(0xF1, 0xF5, 0xFA)
EDGE   = RGBColor(0xD6, 0xE1, 0xEE)
GREEN  = RGBColor(0x0F, 0x6B, 0x45)
AMBER  = RGBColor(0xA8, 0x5A, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

# Usable content band: below the heading, above the bottom accent bar.
X0, X1 = 0.45, 9.55
Y_TOP, Y_BOT = 1.42, 5.34
W = X1 - X0


def txbox(slide, x, y, w, h):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return b, tf


def para(tf, text, size=11, bold=False, color=INK, space_after=4, first=False,
         align=PP_ALIGN.LEFT, italic=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def rich(tf, parts, size=11, space_after=4, first=False, line=None):
    """One paragraph, several differently-styled runs."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return p


def card(slide, x, y, w, h, fill=CARD, edge=EDGE, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if edge is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = edge
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def chip(slide, x, y, w, h, text, fill=NAVY, fg=WHITE, size=9.5, bold=True):
    s = card(slide, x, y, w, h, fill=fill, edge=None, radius=0.5)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size=size, bold=bold, color=fg, space_after=0,
         first=True, align=PP_ALIGN.CENTER)
    return s


def numbered(slide, x, y, n, diameter=0.30, fill=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                               Inches(diameter), Inches(diameter))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, str(n), size=11, bold=True, color=WHITE, space_after=0,
         first=True, align=PP_ALIGN.CENTER)
    return s


def stat(slide, x, y, w, value, label, color=NAVY, vsize=26, lsize=9):
    b, tf = txbox(slide, x, y, w, 0.78)
    para(tf, value, size=vsize, bold=True, color=color, space_after=1, first=True)
    para(tf, label, size=lsize, color=MUTED, space_after=0, line=1.05)
    return b


def heading(slide, text=None, size=19):
    """Restyle the template's own heading box; optionally change its wording."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            if text is not None:
                for p in list(tf.paragraphs)[1:]:
                    p._element.getparent().remove(p._element)
                p0 = tf.paragraphs[0]
                for r in list(p0.runs)[1:]:
                    r._r.getparent().remove(r._r)
                if not p0.runs:
                    p0.add_run()
                p0.runs[0].text = text
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.bold = True
                    r.font.color.rgb = NAVY
                    r.font.name = FONT
            return sh
    return None


def wipe(slide, keep_pictures=True):
    """Remove every text shape, leaving the branded background picture."""
    for sh in list(slide.shapes):
        if sh.shape_type == 13 and keep_pictures:      # PICTURE
            continue
        sh._element.getparent().remove(sh._element)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# =============================================================================
prs = Presentation("template.pptx")
S = prs.slides


# --- 1. COVER (replaces the guidelines page) ---------------------------------
s = S[0]
wipe(s)
b, tf = txbox(s, X0, 1.55, 8.9, 1.5)
para(tf, "Product Intelligence for Industrial Commerce", size=32, bold=True,
     color=NAVY, space_after=6, first=True, line=0.95)
para(tf, "An enrichment engine that refuses to guess.", size=15, color=BLUE,
     bold=True, space_after=0)

b, tf = txbox(s, X0, 3.05, 8.6, 0.62)
para(tf, "A wrong specification ships a broken machine. So every value carries its "
         "evidence, and the engine leaves a field blank, flags the contradiction, or "
         "blocks publication rather than emit a number it cannot defend.",
     size=11.5, color=MUTED, space_after=0, first=True, line=1.25)

for i, (v, l) in enumerate([
        ("2.75×", "attribute coverage lift"),
        ("0.0%", "contradictions on\nevidence-backed values"),
        ("100%", "seeded defects caught\n(51 / 51)"),
        ("$0.0084", "per SKU, batched")]):
    stat(s, X0 + i * 2.28, 3.95, 2.15, v, l, vsize=25)

b, tf = txbox(s, X0, 5.02, 9.0, 0.3)
rich(tf, [("Live prototype  ", True, NAVY),
          ("industrial-product-intelligence.onrender.com", False, BLUE),
          ("     ·     102-case benchmark reproduces from a clean clone", False, MUTED)],
     size=9.5, space_after=0, first=True)
note(s, "One line: minimal input in, commerce-ready record out, and every value "
        "defensible. The four numbers are all reproducible from the repo.")


# --- 2. TEAM DETAILS ---------------------------------------------------------
s = S[1]
wipe(s)
TEAM_NAME = "Team Codes"
TEAM_LEADER = "Ankur Goswami"
TEAM_MEMBERS = ("Ashwin S", "Argho Kumar Halder", "SV Chiranjeevi")

b, tf = txbox(s, 0.55, 3.62, 8.9, 1.6)
para(tf, "Team Details", size=20, bold=True, color=NAVY, space_after=10, first=True)
for label, value in (("Team name", TEAM_NAME),
                     ("Team leader", TEAM_LEADER),
                     ("Members", "  ·  ".join(TEAM_MEMBERS))):
    rich(tf, [(f"{label}:   ", True, INK), (value, False, INK)],
         size=12, space_after=7)
note(s, "Team Codes — leader Ankur Goswami, with Ashwin S, Argho Kumar Halder and "
        "SV Chiranjeevi. Emails are on the portal roster and deliberately not on the "
        "slide; add them here if the organizers ask for them.")


# --- 3. BRIEF ABOUT YOUR SOLUTION -------------------------------------------
s = S[2]
heading(s)
b, tf = txbox(s, X0, 1.5, W, 0.80)
para(tf, "Given a messy catalogue row, produce a complete, standardised, "
         "search-ready product record — with every value traceable.",
     size=13, bold=True, color=INK, space_after=4, first=True, line=1.15)
para(tf, "The engine is deterministic first. AI is a bounded contributor, never the "
         "source of record.", size=10.5, color=MUTED, space_after=0)

stages = ["normalise", "classify", "extract", "infer", "reconcile",
          "vocabulary", "content", "compliance", "validate", "score"]
cw, gap = 0.855, 0.05
for i, st in enumerate(stages):
    x = X0 + i * (cw + gap)
    fill = BLUE if st in ("vocabulary", "compliance") else NAVY
    chip(s, x, 2.36, cw, 0.32, st, fill=fill, size=8)

b, tf = txbox(s, X0, 2.76, W, 0.24)
para(tf, "One pipeline. A value read off page 2 of a datasheet earns the same "
         "provenance, validation and score as a hand-typed one.",
     size=9.5, color=MUTED, space_after=0, first=True, italic=True)

items = [
    ("Every value is attributable",
     "Six provenance classes, ranked. When two stages disagree the stronger wins, "
     "the conflict is recorded, and the winner's confidence is reduced."),
    ("Refusal is a feature, not a failure",
     "Outside the approved vocabulary, not an exact 64th, no manufacturer source — "
     "the engine says so and blocks, instead of guessing plausibly."),
    ("Written the customer's way",
     "The same product rendered five times to five character limits, reproduced "
     "exactly against Unilog's own labelled rows."),
]
cw = (W - 0.4) / 3
for i, (t, d) in enumerate(items):
    x = X0 + i * (cw + 0.2)
    card(s, x, 3.15, cw, 1.62)
    b, tf = txbox(s, x + 0.18, 3.32, cw - 0.36, 1.3)
    para(tf, t, size=11, bold=True, color=NAVY, space_after=5, first=True, line=1.05)
    para(tf, d, size=9.5, color=MUTED, space_after=0, line=1.2)
note(s, "The pipeline is one path for every input. Vocabulary and compliance "
        "(highlighted) are the Unilog-specific stages.")


# --- 4. THE THREE QUESTIONS --------------------------------------------------
s = S[3]
wipe(s)
b, tf = txbox(s, X0, 0.85, W, 0.4)
para(tf, "How the solution enriches, proves and scales", size=19, bold=True,
     color=NAVY, space_after=0, first=True)

blocks = [
    ("1", "Enriching minimal input",
     "Part number and brand decode against ISO 15 / ISO 898-1 and a curated knowledge "
     "base; datasheets, product pages and catalogue rows all feed the same pipeline.",
     [("2.75×", "coverage lift"), ("61.1%", "withheld values recovered")]),
    ("2", "Accuracy and trust",
     "16 cross-field engineering rules catch contradictions that are individually "
     "plausible. Readiness scores on three axes and returns publish / review / blocked.",
     [("100%", "defects caught (51/51)"), ("0.0%", "false alarms on clean records")]),
    ("3", "Scale",
     "611 rows/s ingested and 305 products/s enriched on one core. Unseen categories "
     "are learned from the data, not hand-curated.",
     # The condition travels with the number: 81.7% is what a reviewer who
     # approves only well-supported clusters gets. Approving all 83 proposals
     # reaches 91.4%. Quoting either alone is the 14/14 mistake again.
     [("8.9% → 81.7%", "their 1,000 rows — 38 of\n83 clusters approved"),
      ("$0.0084", "per SKU, batched")]),
]
y = 1.42
for n, title, body, stats in blocks:
    card(s, X0, y, W, 1.20)
    numbered(s, X0 + 0.22, y + 0.20, n)
    b, tf = txbox(s, X0 + 0.66, y + 0.17, 5.55, 0.95)
    para(tf, title, size=12, bold=True, color=NAVY, space_after=4, first=True)
    para(tf, body, size=9.5, color=MUTED, space_after=0, line=1.2)
    for j, (v, l) in enumerate(stats):
        bx, btf = txbox(s, X0 + 6.35 + j * 1.35, y + 0.30, 1.3, 0.6)
        para(btf, v, size=13, bold=True, color=BLUE, space_after=1, first=True)
        para(btf, l, size=7.5, color=MUTED, space_after=0, line=1.05)
    y += 1.30
note(s, "Three answers, three sets of measured evidence. Nothing here is an estimate "
        "except the batched cost, which is labelled as a projection in the repo.")


# --- 5. OPPORTUNITIES / USP --------------------------------------------------
s = S[4]
heading(s, "Opportunities — and what makes this different")
b, tf = txbox(s, X0, 1.45, W, 0.5)
rich(tf, [("Coverage is cheap. ", True, NAVY),
          ("Any system reaches 100% by inventing everything. The hard part — and the "
           "only part a distributor can actually publish — is precision you can defend.",
           False, INK)], size=12, space_after=0, first=True, line=1.2)

card(s, X0, 2.10, 4.42, 1.42)
b, tf = txbox(s, X0 + 0.2, 2.26, 4.05, 1.1)
para(tf, "The USP", size=11, bold=True, color=NAVY, space_after=5, first=True)
para(tf, "It is the only enrichment engine we know of that reports what it "
         "refused to do. Every AI proposal, vocabulary mismatch and rejected "
         "source is listed with its reason — on screen and in the export.",
     size=9.5, color=MUTED, space_after=0, line=1.22)

card(s, X0 + 4.68, 2.10, 4.42, 1.42)
b, tf = txbox(s, X0 + 4.88, 2.26, 4.05, 1.1)
para(tf, "Proven, not asserted", size=11, bold=True, color=NAVY, space_after=5, first=True)
para(tf, "A live ablation measured the AI making accuracy worse; bounding it to "
         "gap-filling beat both the raw model and the engine alone. An independent "
         "A–Z pass then found five defects — all fixed, each with a regression, "
         "every benchmark figure unchanged.",
     size=9.5, color=MUTED, space_after=0, line=1.18)

b, tf = txbox(s, X0, 3.72, W, 0.3)
para(tf, "Measured against Unilog's own labelled delivery rows", size=10.5,
     bold=True, color=NAVY, space_after=0, first=True)
for i, (v, l, c) in enumerate([
        ("14 / 14", "prose fields exact — given\nthe attribute values", GREEN),
        ("252", "delivery columns emitted\nin their exact order", NAVY),
        ("61.5%", "precision on the fields\nwe do answer", NAVY),
        ("88 of 114", "left blank: the data is not\nin the input row", AMBER)]):
    stat(s, X0 + i * 2.28, 4.08, 2.2, v, l, color=c, vsize=19, lsize=8.5)
note(s, "Say the condition out loud on 14/14: it is the formatter given the "
        "values, not the pipeline from a catalogue row — slide 8 shows both. The "
        "last number is the deliberate one: three quarters of their format cannot "
        "be reached from a 6-column row at all.")


# --- 6. FEATURES -------------------------------------------------------------
s = S[5]
heading(s)
feats = [
    ("Provenance on every value", "Six ranked classes from supplied to defaulted, with evidence text and a source URL."),
    ("Cross-field validation", "16 rules, plus a part number's own standard challenging a supplied value it contradicts."),
    ("Bounded AI gate", "The model may fill a gap or displace a default. It may never overrule evidence."),
    ("Unilog compliance layer", "Approved UOM, exact-64th fractions, list-of-values refusals, character limits — on screen."),
    ("Auto-taxonomy learning", "Clusters unclassified rows and proposes a schema; a human approves it."),
    ("Manufacturer-only sourcing", "Marketplaces, retailers and distributors refused — each with a stated reason."),
]
cw = (W - 0.22) / 2
ch = 0.85
for i, (t, d) in enumerate(feats):
    x = X0 + (i % 2) * (cw + 0.22)
    y = 1.50 + (i // 2) * (ch + 0.17)
    card(s, x, y, cw, ch)
    numbered(s, x + 0.20, y + 0.26, i + 1, diameter=0.28,
             fill=NAVY if i % 2 == 0 else BLUE)
    b, tf = txbox(s, x + 0.62, y + 0.15, cw - 0.82, ch - 0.24)
    para(tf, t, size=11, bold=True, color=NAVY, space_after=3, first=True)
    para(tf, d, size=9, color=MUTED, space_after=0, line=1.18)
note(s, "Six shipped capabilities, each with a test that proves the specific "
        "behaviour it claims and costs nothing to run.")


# --- 7. PROCESS FLOW ---------------------------------------------------------
s = S[6]
heading(s)
b, tf = txbox(s, X0, 1.42, W, 0.26)
para(tf, "Four input paths converge on one RawProduct, then one pipeline. "
         "Refusal points are marked in amber.",
     size=10, color=MUTED, space_after=0, first=True)

inputs = ["Form / CSV row", "Unilog catalogue row", "Datasheet PDF", "Brand + part number"]
for i, t in enumerate(inputs):
    y = 1.82 + i * 0.44
    card(s, X0, y, 2.05, 0.36, fill=WHITE, edge=EDGE)
    b, tf = txbox(s, X0 + 0.12, y + 0.08, 1.85, 0.22)
    para(tf, t, size=8.5, bold=True, color=NAVY, space_after=0, first=True)

card(s, 2.75, 1.82, 1.32, 1.78, fill=NAVY, edge=None)
b, tf = txbox(s, 2.85, 2.44, 1.12, 0.6)
para(tf, "RawProduct", size=10, bold=True, color=WHITE, space_after=2, first=True,
     align=PP_ALIGN.CENTER)
para(tf, "one shape", size=8, color=RGBColor(0xBF, 0xD8, 0xF2), space_after=0,
     first=False, align=PP_ALIGN.CENTER)

steps = [("classify", "category + schema"), ("extract", "spec table & prose"),
         ("gate", "AI may add, never overrule"), ("vocabulary", "approved values only"),
         ("compliance", "house style & limits"), ("score", "publish / review / blocked")]
for i, (t, d) in enumerate(steps):
    x = 4.32 + (i % 3) * 1.78
    y = 1.82 + (i // 3) * 0.92
    flagged = t in ("gate", "vocabulary")
    card(s, x, y, 1.66, 0.78, fill=CARD, edge=EDGE)
    b, tf = txbox(s, x + 0.13, y + 0.13, 1.4, 0.55)
    para(tf, t, size=10, bold=True, color=AMBER if flagged else NAVY,
         space_after=2, first=True)
    para(tf, d, size=7.5, color=MUTED, space_after=0, line=1.12)

card(s, X0, 3.82, W, 0.74, fill=WHITE, edge=EDGE)
b, tf = txbox(s, X0 + 0.2, 3.95, W - 0.4, 0.52)
rich(tf, [("Output   ", True, NAVY),
          ("EnrichedProduct → 252-column Unilog delivery sheet · schema.org JSON-LD · "
           "catalogue CSV. Every attribute carries provenance, confidence, evidence "
           "and the URL it was read from.", False, MUTED)],
     size=9.5, space_after=0, first=True, line=1.2)

b, tf = txbox(s, X0, 4.72, W, 0.5)
rich(tf, [("At each amber point the engine can stop. ", True, AMBER),
          ("Across the 102-case benchmark the gate refused 28 of the model's 95 "
           "proposals; a value outside the approved list keeps its original and raises "
           "a violation rather than being rewritten to the nearest match.", False, INK)],
     size=9.5, space_after=0, first=True, line=1.2)
note(s, "The amber boxes are where the engine is allowed to say no. That is the "
        "whole design in one diagram.")


# --- 8. EXPECTED OUTPUT vs OURS ----------------------------------------------
# Built from run_expected_vs_ours.py so the strings on the slide are the ones
# the scorer actually compared, and cannot drift from the measurement.
import json as _json
_gt = _json.load(open("expected_vs_ours.json", encoding="utf-8"))
_row = _gt["cases"][0]
_sum = _gt["summary"]

s = S[7]
heading(s, "Their expected output, next to ours")
b, tf = txbox(s, X0, 1.40, W, 0.44)
rich(tf, [("Two fully worked delivery rows — the only labelled ground truth "
           "in their pack. ", True, NAVY),
          ("Row 1 below, their value then ours. Every prose field is built by "
           "formula against their character limits, never free-hand.",
           False, MUTED)],
     size=10, space_after=0, first=True, line=1.2)

# header strip
hdr = [("FIELD", X0 + 0.02, 1.05), ("UNILOG EXPECTED", X0 + 1.15, 5.35),
       ("OURS", X0 + 6.60, 2.45)]
for label, x, w in hdr:
    b, tf = txbox(s, x, 1.94, w, 0.18)
    para(tf, label, size=7.5, bold=True, color=MUTED, space_after=0, first=True)

y = 2.16
for f in _row["fields"][:4]:
    h = 0.56 if len(f["expected"]) > 78 else 0.42
    card(s, X0, y, W, h, fill=CARD)
    b, tf = txbox(s, X0 + 0.14, y + 0.09, 1.02, 0.3)
    para(tf, f["column"].replace("_DESC", ""), size=8, bold=True, color=NAVY,
         space_after=0, first=True)
    b, tf = txbox(s, X0 + 1.15, y + 0.08, 5.35, h - 0.16)
    para(tf, f["expected"], size=7.5, color=INK, space_after=0, first=True, line=1.15)
    b, tf = txbox(s, X0 + 6.60, y + 0.09, 2.45, 0.3)
    if f["match_given"]:
        rich(tf, [("identical", True, GREEN),
                  ("  — character for character", False, MUTED)],
             size=7.5, space_after=0, first=True)
    else:
        para(tf, f["given_attributes"][:40], size=7.5, color=AMBER,
             space_after=0, first=True)
    y += h + 0.06

# The condition, stated on the slide rather than in a footnote.
card(s, X0, y + 0.04, W, 0.86, fill=WHITE, edge=EDGE)
b, tf = txbox(s, X0 + 0.18, y + 0.14, W - 0.36, 0.68)
rich(tf, [(f"{_sum['exact_given_attributes']}/{_sum['fields_scored']} exact when the "
           f"attribute values are supplied. ", True, GREEN),
          (f"{_sum['exact_from_input_row']}/{_sum['fields_scored']} from the six-column "
           f"input row alone", True, AMBER),
          (" — one engine, two different inputs. Series, Mounting, Wash Cycles and "
           "Voltage are not in a catalogue row; they are on frigidaire.com, which "
           "their own delivery row cites as MFR URL. The formatting is solved. The "
           "sourcing is not, and 88 of 114 blank fields is the size of it.", False, INK)],
     size=9, space_after=0, first=True, line=1.22)
note(s, "Both numbers appear only here, so say both out loud. 14/14 alone claims the "
        "pipeline produces these strings from a catalogue row, and it does not — it "
        "says the formatter is exact once it has the values. The distance between the "
        "two is the sourcing gap, and it is the honest headline of this slide.")


# --- 9. ARCHITECTURE ---------------------------------------------------------
s = S[8]
heading(s)
tiers = [
    ("Ingest", BLUE, ["form / CSV", "Unilog rows", "PDF datasheet", "product page", "discovery"]),
    ("Pipeline", NAVY, ["normalise", "classify", "extract", "infer + gate",
                        "reconcile", "vocabulary", "content", "compliance",
                        "validate", "score"]),
    ("Output", BLUE, ["Unilog 252-col", "schema.org", "catalogue CSV", "refusal ledger"]),
]
y = 1.50
for name, colour, items in tiers:
    h = 0.92 if name == "Pipeline" else 0.62
    card(s, X0, y, W, h, fill=WHITE, edge=EDGE)
    b, tf = txbox(s, X0 + 0.16, y + 0.10, 1.15, 0.4)
    para(tf, name, size=10.5, bold=True, color=colour, space_after=0, first=True)
    per_row = 5
    cw = 1.42
    for i, it in enumerate(items):
        cx = X0 + 1.42 + (i % per_row) * (cw + 0.09)
        cy = y + 0.12 + (i // per_row) * 0.38
        chip(s, cx, cy, cw, 0.30, it, fill=CARD, fg=NAVY, size=8, bold=False)
    y += h + 0.16

card(s, X0, y, 4.42, 1.06, fill=CARD)
b, tf = txbox(s, X0 + 0.18, y + 0.14, 4.05, 0.8)
para(tf, "Providers are swappable", size=10.5, bold=True, color=NAVY,
     space_after=4, first=True)
para(tf, "Deterministic engine (free, default) · Claude tool-use (live) · Hybrid, "
         "where the model may only fill a gap or displace a default.",
     size=9, color=MUTED, space_after=0, line=1.2)

card(s, X0 + 4.68, y, 4.42, 1.06, fill=CARD)
b, tf = txbox(s, X0 + 4.86, y + 0.14, 4.05, 0.8)
para(tf, "The customer's rules are data, not code", size=10.5, bold=True,
     color=NAVY, space_after=4, first=True)
para(tf, "Taxonomy, UOM table, list of values, export schema and source policy are "
         "JSON. Their spreadsheets land as a data drop, not a rewrite.",
     size=9, color=MUTED, space_after=0, line=1.2)
note(s, "One pipeline, swappable providers, and every customer-specific rule held "
        "as data so onboarding a new schema is a file, not a sprint.")


# --- 10. TECHNOLOGIES --------------------------------------------------------
s = S[9]
heading(s)
groups = [
    ("Engine", ["Python 3.13", "FastAPI", "Pydantic v2", "pdfplumber", "BeautifulSoup + lxml", "httpx"]),
    ("AI", ["Claude Sonnet", "tool-schema structured output", "server-side web search", "content-addressed cache"]),
    ("Interface", ["React 19", "Vite", "zero runtime UI dependencies", "design tokens, dark mode"]),
    ("Delivery", ["Docker (single service)", "Render", "UptimeRobot", "GitHub"]),
]
cw = (W - 0.33) / 2
for i, (name, items) in enumerate(groups):
    x = X0 + (i % 2) * (cw + 0.33)
    y = 1.50 + (i // 2) * 1.82
    b, tf = txbox(s, x, y, cw, 0.28)
    para(tf, name, size=11, bold=True, color=NAVY, space_after=0, first=True)
    for j, it in enumerate(items):
        cx = x + (j % 2) * (cw / 2)
        cy = y + 0.36 + (j // 2) * 0.38
        chip(s, cx, cy, cw / 2 - 0.12, 0.30, it, fill=CARD, fg=INK, size=8.5, bold=False)

b, tf = txbox(s, X0, 5.02, W, 0.3)
rich(tf, [("Deliberate constraint  ", True, NAVY),
          ("the deployed instance holds no API key and cannot spend. Reviewers see "
           "genuine model output from 20 pre-computed records at zero cost.",
           False, MUTED)], size=9, space_after=0, first=True)
note(s, "Boring, current, and cheap to run. The interesting choice is the one at the "
        "bottom: the public demo is structurally unable to cost anyone money.")


# --- 11. COST ----------------------------------------------------------------
s = S[10]
heading(s, "Cost per SKU — measured, then projected")
b, tf = txbox(s, X0, 1.44, W, 0.3)
para(tf, "Their baseline: 10 minutes of analyst time per SKU at $35/hr.",
     size=10.5, color=MUTED, space_after=0, first=True)

rows = [
    ("Manual, today", "$5.83", "per SKU", "$4,372,500", MUTED, False),
    ("Every SKU to the model", "$0.02381", "standard rates", "$17,857", INK, False),
    ("+ deterministic-first triage", "$0.01681", "70.6% call rate", "$12,605", INK, False),
    ("+ Batch API", "$0.00840", "projected", "$6,302", NAVY, True),
]
y = 1.86
for label, per, note_txt, monthly, colour, hero in rows:
    card(s, X0, y, W, 0.62, fill=CARD if hero else WHITE,
         edge=None if hero else EDGE)
    b, tf = txbox(s, X0 + 0.2, y + 0.19, 3.5, 0.3)
    para(tf, label, size=10.5, bold=hero, color=NAVY if hero else INK,
         space_after=0, first=True)
    b, tf = txbox(s, X0 + 3.9, y + 0.13, 1.7, 0.42)
    para(tf, per, size=14 if hero else 12, bold=True, color=colour,
         space_after=0, first=True)
    b, tf = txbox(s, X0 + 5.7, y + 0.21, 1.6, 0.26)
    para(tf, note_txt, size=8.5, color=MUTED, space_after=0, first=True)
    b, tf = txbox(s, X0 + 7.3, y + 0.17, 1.6, 0.32)
    para(tf, monthly, size=11.5, bold=True, color=colour, space_after=0,
         first=True, align=PP_ALIGN.RIGHT)
    y += 0.70

b, tf = txbox(s, X0, 4.72, 5.6, 0.5)
rich(tf, [("0.14% of manual cost", True, GREEN),
          ("  at 750,000 SKUs/month — the organizers' own scaling target.",
           False, INK)], size=11, space_after=0, first=True)
b, tf = txbox(s, X0 + 5.8, 4.70, 3.3, 0.56)
para(tf, "Binding constraint is the API rate limit, not compute: 0.7 compute "
         "hours/month, 0.20 calls/second sustained.",
     size=8.5, color=MUTED, space_after=0, first=True, line=1.2)
note(s, "Right-hand column is monthly cost at 750k SKUs. Triage and token counts are "
        "measured; the batch figure is a projection from published rates and is "
        "labelled as such in the repo.")


# --- 12. MVP SNAPSHOTS -------------------------------------------------------
# Captured from the deployed prototype by docs/deck/shots.py, cropped to 2.7:1
# so four fit the slide. Each keeps the top of its card, where the headline
# badges and numbers are; the detail rows would be unreadable at deck scale.
s = S[11]
heading(s)
# Rides on the heading's own line, right-aligned: stacked underneath it
# collided with the 19pt descenders, and there is no vertical room to spare.
b, tf = txbox(s, 5.30, 1.00, 4.25, 0.22)
para(tf, "Captured from the live prototype — the same link a judge opens.",
     size=9, color=MUTED, space_after=0, first=True, italic=True,
     align=PP_ALIGN.RIGHT)

shots = [
    ("shots/crop/1_enrichment.png", "Single product enrichment",
     "Readiness on three axes, category with its reasoning, provenance per value"),
    ("shots/crop/2_gate.png", "The AI gate refusing",
     "Two proposals refused against ISO 15; one unbacked default replaced"),
    ("shots/crop/3_content.png", "The customer's content standard",
     "Five formats to five limits; the 40-character line names the facts it dropped"),
    ("shots/crop/4_catalog.png", "Catalog at volume",
     "10 products in under a second — 7 publishable, 1 review, 2 blocked"),
]
cw = (W - 0.28) / 2
ih = cw / 2.70
for i, (path, title, desc) in enumerate(shots):
    x = X0 + (i % 2) * (cw + 0.28)
    y = 1.42 + (i // 2) * 2.05
    pic = s.shapes.add_picture(path, Inches(x), Inches(y),
                               width=Inches(cw), height=Inches(ih))
    pic.line.color.rgb = EDGE
    pic.line.width = Pt(0.75)
    b, tf = txbox(s, x, y + ih + 0.03, cw, 0.33)
    para(tf, title, size=9.5, bold=True, color=NAVY, space_after=1, first=True)
    para(tf, desc, size=8, color=MUTED, space_after=0, line=1.1)
note(s, "All four are the deployed app, not mockups. The gate panel is the one to "
        "talk to: the model proposed 14.8 kN and 14000 rpm, ISO 15 says 14 kN and "
        "16000, and both were refused while an unbacked default was replaced.")


# --- 13. FUTURE --------------------------------------------------------------
s = S[12]
heading(s, "What we measured but did not solve")
b, tf = txbox(s, X0, 1.44, W, 0.3)
para(tf, "Each item below is a known gap with a number attached, not a wish list.",
     size=10.5, color=MUTED, space_after=0, first=True)

future = [
    ("Headless rendering for manufacturer sites", AMBER,
     "Zero of four real sites yielded a spec table: Frigidaire, Milwaukee and SKF all "
     "return HTTP 200 with client-side-rendered content. This is the single change that "
     "moves delivery accuracy past 14%."),
    ("The seven reference files still to land", NAVY,
     "UOM standards, content guidelines, the 161k-row list of values and the 27k-row "
     "brand list. Every stub reports source: provisional, and swapping them in is a "
     "data drop — the seams are already built."),
    ("Attribute recall on categories without a standard", NAVY,
     "41.7% against 99.5% where ISO fixes the answer. Honest ceiling of a knowledge "
     "base; closing it needs manufacturer documents, not a better prompt."),
    ("Batch API and video extraction", MUTED,
     "Costed at $0.0084/SKU but not implemented, and flagged by the organizers as "
     "innovative but slow. Both are catalogue-path work, not interactive."),
]
y = 1.82
for t, colour, d in future:
    card(s, X0, y, W, 0.80)
    b, tf = txbox(s, X0 + 0.2, y + 0.12, W - 0.4, 0.6)
    para(tf, t, size=10.5, bold=True, color=colour, space_after=3, first=True)
    para(tf, d, size=9, color=MUTED, space_after=0, line=1.18)
    y += 0.88
note(s, "Saying what does not work, with the measurement, is the same discipline the "
        "engine applies to a product record.")


# --- 14. LINKS ---------------------------------------------------------------
s = S[13]
wipe(s)
b, tf = txbox(s, X0, 0.88, W, 0.4)
para(tf, "Links", size=19, bold=True, color=NAVY, space_after=0, first=True)

links = [
    ("Working prototype", "https://industrial-product-intelligence.onrender.com",
     "Live, monitored, no API key needed — every tab is explorable at zero cost."),
    ("GitHub repository", "https://github.com/ashwinsureshh/industrial-product-intelligence",
     "Public. 9 test suites, 337 assertions, all free to run. Benchmarks reproduce from a clean clone."),
    ("Demo video (3 min)", "[ paste link before submitting ]",
     "Walkthrough of enrichment, the AI gate refusing two values, and the delivery export."),
]
y = 1.55
for t, url, d in links:
    card(s, X0, y, W, 1.06)
    b, tf = txbox(s, X0 + 0.22, y + 0.16, W - 0.44, 0.8)
    para(tf, t, size=11, bold=True, color=NAVY, space_after=4, first=True)
    para(tf, url, size=10.5, bold=True, color=BLUE, space_after=4)
    para(tf, d, size=9, color=MUTED, space_after=0, line=1.18)
    y += 1.20

b, tf = txbox(s, X0, 5.02, W, 0.3)
para(tf, "Reproduce the headline: cd backend && python run_benchmark.py    "
         "(102 cases, no API key, $0)",
     size=9, color=MUTED, space_after=0, first=True, italic=True)
note(s, "PASTE the demo video link before submitting, and confirm the repo is public.")

prs.save("UniHack_Prototype_Submission.pptx")
print("saved UniHack_Prototype_Submission.pptx")

# The submitted deliverable lives one directory up. Copying it here rather than
# leaving it as a remembered step: a rebuild that updates only the working copy
# looks like it worked and ships the previous deck.
import shutil as _shutil
from pathlib import Path as _Path
_deliverable = _Path(__file__).resolve().parent.parent / "UniHack_Prototype_Submission.pptx"
_shutil.copyfile("UniHack_Prototype_Submission.pptx", _deliverable)
print(f"copied to {_deliverable}")
