"""Analytical overflow check: measure every text box with real Arial metrics.

No LibreOffice in this environment, so slides cannot be rendered for visual QA.
This substitutes for it on the defect that actually matters — text overflowing
its shape — by laying out each paragraph with PIL against the real font file and
comparing the wrapped height to the shape height.

Reports both directions: text taller than its box (clipping) and text spilling
past the slide edge.
"""
import sys
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
SLIDE_W, SLIDE_H = 10.0, 5.625

FONTS = {
    (False, False): r"C:\Windows\Fonts\arial.ttf",
    (True, False):  r"C:\Windows\Fonts\arialbd.ttf",
    (False, True):  r"C:\Windows\Fonts\ariali.ttf",
    (True, True):   r"C:\Windows\Fonts\arialbi.ttf",
}
_cache = {}


def font(size_pt, bold, italic):
    key = (round(size_pt * 2), bold, italic)
    if key not in _cache:
        # PIL sizes in px; at 96 dpi 1pt = 1.333px
        _cache[key] = ImageFont.truetype(FONTS[(bold, italic)], int(size_pt * 96 / 72))
    return _cache[key]


def wrap_lines(text, f, max_px):
    if not text:
        return 1
    lines, cur = 0, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if f.getlength(trial) <= max_px or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines + (1 if cur else 0)


def check(path):
    prs = Presentation(path)
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            text_all = tf.text.strip()
            if not text_all:
                continue

            left, top = sh.left / EMU, sh.top / EMU
            w, h = sh.width / EMU, sh.height / EMU
            ml = (tf.margin_left or 0) / EMU
            mr = (tf.margin_right or 0) / EMU
            mt = (tf.margin_top or 0) / EMU
            mb = (tf.margin_bottom or 0) / EMU
            avail_px = max((w - ml - mr) * 96, 1)

            total_h = 0.0
            for p in tf.paragraphs:
                runs = p.runs
                if not runs:
                    continue
                size = max((r.font.size.pt if r.font.size else 11) for r in runs)
                bold = any(bool(r.font.bold) for r in runs)
                ital = any(bool(r.font.italic) for r in runs)
                ptext = "".join(r.text for r in runs)
                f = font(size, bold, ital)
                nlines = 0
                for seg in ptext.split("\n"):
                    nlines += wrap_lines(seg, f, avail_px)
                ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
                total_h += nlines * size * 1.21 * ls / 72.0
                total_h += (p.space_after.pt / 72.0) if p.space_after else 0.0

            box_h = h - mt - mb
            if total_h > box_h + 0.02:
                problems.append(
                    f"s{idx:>2} OVERFLOW  {total_h:.2f}in of text in a {box_h:.2f}in box"
                    f"  :: {text_all[:56]!r}")
            # Template-owned boxes keep their original generous geometry and are
            # nearly empty; only flag margins on shapes this build added.
            if not sh.name.startswith("Google Shape"):
                if (left < 0.28 or top < 0.05
                        or left + w > SLIDE_W - 0.28 + 0.01
                        or top + h > SLIDE_H - 0.16 + 0.01):
                    problems.append(
                        f"s{idx:>2} MARGIN    L={left:.2f} T={top:.2f} R={left+w:.2f} "
                        f"B={top+h:.2f}  :: {text_all[:44]!r}")
        # Overlap between shapes this build added (the template background and
        # its near-empty heading boxes are excluded).
        mine = [sh for sh in slide.shapes
                if not sh.name.startswith("Google Shape") and sh.shape_type != 13]
        for i, a in enumerate(mine):
            for b_ in mine[i + 1:]:
                ax0, ay0 = a.left / EMU, a.top / EMU
                ax1, ay1 = ax0 + a.width / EMU, ay0 + a.height / EMU
                bx0, by0 = b_.left / EMU, b_.top / EMU
                bx1, by1 = bx0 + b_.width / EMU, by0 + b_.height / EMU
                ox = min(ax1, bx1) - max(ax0, bx0)
                oy = min(ay1, by1) - max(ay0, by0)
                if ox > 0.02 and oy > 0.02:
                    a_txt = a.text_frame.text.strip()[:22] if a.has_text_frame else a.shape_type
                    b_txt = b_.text_frame.text.strip()[:22] if b_.has_text_frame else b_.shape_type
                    # A text box sitting inside its own card is intended.
                    contained = (ax0 >= bx0 - 0.01 and ax1 <= bx1 + 0.01
                                 and ay0 >= by0 - 0.01 and ay1 <= by1 + 0.01) or                                 (bx0 >= ax0 - 0.01 and bx1 <= ax1 + 0.01
                                 and by0 >= ay0 - 0.01 and by1 <= ay1 + 0.01)
                    if not contained:
                        problems.append(
                            f"s{idx:>2} OVERLAP   {ox:.2f}x{oy:.2f}in "
                            f":: {a_txt!r} / {b_txt!r}")
    return problems


if __name__ == "__main__":
    probs = check(sys.argv[1] if len(sys.argv) > 1 else "UniHack_Prototype_Submission.pptx")
    if not probs:
        print("no overflow or margin problems found")
    else:
        print(f"{len(probs)} problem(s):")
        for p in probs:
            print("  " + p)
