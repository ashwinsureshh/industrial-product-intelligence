"""Approximate render of one slide, for eyeballing placement.

Not a PowerPoint renderer: it re-lays the deck's own geometry in HTML so
crowding, collisions and alignment are visible. Font substitution means text
widths are indicative, not exact.
"""
import base64, pathlib, sys
from pptx import Presentation
from playwright.sync_api import sync_playwright

P, SCALE = 914400, 130
# PP_ALIGN: LEFT=1, CENTER=2, RIGHT=3, JUSTIFY=4
ALIGN = {2: "center", 3: "right", 4: "justify"}

idx = int(sys.argv[1]) if len(sys.argv) > 1 else 12
out = sys.argv[2] if len(sys.argv) > 2 else f"slide{idx}_preview.png"
s = Presentation("UniHack_Prototype_Submission.pptx").slides[idx - 1]

parts = []
for sh in s.shapes:
    x, y = sh.left / P * SCALE, sh.top / P * SCALE
    w, h = sh.width / P * SCALE, sh.height / P * SCALE
    if sh.shape_type == 13:
        blob = base64.b64encode(sh.image.blob).decode()
        parts.append(f'<img src="data:image/png;base64,{blob}" style="position:absolute;'
                     f'left:{x}px;top:{y}px;width:{w}px;height:{h}px;object-fit:fill">')
        continue
    fill = ""
    try:
        if sh.fill.type is not None and sh.fill.type == 1:
            fill = f"background:#{sh.fill.fore_color.rgb};"
    except Exception:
        pass
    radius = "border-radius:7px;" if "Rounded" in str(sh.shape_type) else ""
    body = ""
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            al = ALIGN.get(getattr(para.alignment, "value", None), "left")
            runs = ""
            for r in para.runs:
                sz = (r.font.size.pt if r.font.size else 11) * SCALE / 72
                try:
                    col = f"#{r.font.color.rgb}"
                except Exception:
                    col = "#16202c"
                runs += (f'<span style="font-size:{sz:.1f}px;font-weight:'
                         f'{700 if r.font.bold else 400};color:{col};'
                         f'font-style:{"italic" if r.font.italic else "normal"}">'
                         f'{r.text}</span>')
            sa = (para.space_after.pt if para.space_after else 0) * SCALE / 72
            body += f'<div style="text-align:{al};margin-bottom:{sa:.1f}px">{runs}</div>'
    va = "center" if getattr(sh.text_frame, "vertical_anchor", None) == 3 else "flex-start"
    parts.append(f'<div style="position:absolute;left:{x}px;top:{y}px;width:{w}px;'
                 f'height:{h}px;{fill}{radius}font-family:Arial;line-height:1.2;'
                 f'display:flex;flex-direction:column;justify-content:{va};'
                 f'box-sizing:border-box">{body}</div>')

html = (f'<body style="margin:0"><div style="position:relative;width:{10*SCALE}px;'
        f'height:{5.625*SCALE}px;background:#fff;overflow:hidden">'
        + "".join(parts) + "</div></body>")
pathlib.Path("_preview.html").write_text(html, encoding="utf-8")

with sync_playwright() as p:
    b = p.chromium.launch(channel="chrome", headless=True)
    pg = b.new_page(viewport={"width": 10 * SCALE, "height": int(5.625 * SCALE)},
                    device_scale_factor=2)
    pg.goto(pathlib.Path("_preview.html").resolve().as_uri())
    pg.wait_for_timeout(600)
    pg.screenshot(path=out)
    b.close()
print("wrote", out)
